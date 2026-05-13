# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

# ── 색상 ──────────────────────────────────────────────────────
C = {
    'blue_dark'  : RGBColor(0x1e,0x40,0xaf),
    'blue_mid'   : RGBColor(0x25,0x63,0xeb),
    'blue_light' : RGBColor(0xdb,0xea,0xfe),
    'blue_pale'  : RGBColor(0xef,0xf6,0xff),
    'gray_dark'  : RGBColor(0x11,0x18,0x27),
    'gray_mid'   : RGBColor(0x37,0x41,0x51),
    'gray_light' : RGBColor(0xf3,0xf4,0xf6),
    'white'      : RGBColor(0xff,0xff,0xff),
    'green'      : RGBColor(0x16,0xa3,0x4a),
    'green_light': RGBColor(0xdc,0xfc,0xe7),
    'amber'      : RGBColor(0xd9,0x77,0x06),
    'amber_light': RGBColor(0xfe,0xf3,0xc7),
    'red'        : RGBColor(0xdc,0x26,0x26),
    'red_light'  : RGBColor(0xfe,0xe2,0xe2),
    'purple'     : RGBColor(0x7c,0x3a,0xed),
    'purple_light':RGBColor(0xed,0xe9,0xfe),
    'teal'       : RGBColor(0x06,0xb6,0xd4),
    'teal_light' : RGBColor(0xcc,0xfb,0xf1),
}

def rect(slide, x, y, w, h, fill=None, line_color=None, line_width=1.2, radius=False):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size=13, bold=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or C['gray_dark']
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb

def node(slide, label, sub, x, y, w, h,
         fill=None, border=None, lbl_size=13, sub_size=11):
    fill   = fill   or C['white']
    border = border or C['blue_light']
    rect(slide, x, y, w, h, fill=fill, line_color=border, line_width=1.5)
    mid_y = y + (h - Inches(0.38)) / 2 - Inches(0.05)
    txt(slide, label, x+Inches(0.1), mid_y, w-Inches(0.2), Inches(0.38),
        size=lbl_size, bold=True, color=C['gray_dark'], align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, x+Inches(0.1), mid_y+Inches(0.35), w-Inches(0.2), Inches(0.32),
            size=sub_size, color=C['gray_mid'], align=PP_ALIGN.CENTER)

def arrow_h(slide, x1, y, x2, color=None):
    color = color or C['gray_mid']
    line = slide.shapes.add_connector(1, x1, y, x2, y)
    line.line.color.rgb = color; line.line.width = Pt(2)

def arrow_v(slide, x, y1, y2, color=None):
    color = color or C['gray_mid']
    line = slide.shapes.add_connector(1, x, y1, x, y2)
    line.line.color.rgb = color; line.line.width = Pt(2)

def badge(slide, text, x, y, w, h, fill, fc=None):
    fc = fc or C['white']
    rect(slide, x, y, w, h, fill=fill)
    txt(slide, text, x, y, w, h, size=10, bold=True,
        color=fc, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
#  PPT 1 — 시스템 동작 다이어그램
# ════════════════════════════════════════════════════════════════════
prs1 = Presentation()
prs1.slide_width = W; prs1.slide_height = H
blank = prs1.slide_layouts[6]

sl = prs1.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C['gray_light'])

# ── 헤더 ──────────────────────────────────────────────────────
rect(sl, 0, 0, W, Inches(0.85), fill=C['blue_dark'])
txt(sl, 'Aegis QA Assistant — 시스템 동작 다이어그램',
    Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
    size=22, bold=True, color=C['white'])
txt(sl, 'PDF 업로드 → 메뉴트리 추출 → 사용자 검토 → TC 생성 → 검토 → Excel',
    Inches(0.4), Inches(0.62), Inches(12), Inches(0.28),
    size=12, color=C['blue_light'])

# ── 3개 레인 배경 ──────────────────────────────────────────────
lane_y  = [Inches(0.92), Inches(2.62), Inches(4.62)]
lane_h  = [Inches(1.62), Inches(1.92), Inches(2.62)]
lane_labels = ['사용자 (QA 담당자)', 'React 프론트엔드', 'FastAPI 백엔드 + AI']
lane_fills  = [C['amber_light'], C['blue_pale'], C['green_light']]
lane_lc     = [C['amber'], C['blue_mid'], C['green']]

for i in range(3):
    rect(sl, 0, lane_y[i], W, lane_h[i], fill=lane_fills[i], line_color=lane_lc[i], line_width=0.8)
    txt(sl, lane_labels[i],
        Inches(0.08), lane_y[i]+Inches(0.08),
        Inches(1.5), Inches(0.35),
        size=10, bold=True, color=lane_lc[i])

# ── 사용자 레인 박스 ──────────────────────────────────────────
u_nodes = [
    ('① PDF\n업로드', '기획서 PDF'),
    ('③ 메뉴트리\n검토·편집', '범위 확정'),
    ('⑤ TC 검토\n승인·수정', '품질 확인'),
    ('⑦ Excel\n다운로드', 'TC 리포트'),
]
uw = Inches(1.55); uh = Inches(1.25); uy = Inches(1.0)
ux_starts = [Inches(1.7), Inches(4.4), Inches(7.55), Inches(11.05)]
for i, (lbl, sub) in enumerate(u_nodes):
    node(sl, lbl, sub, ux_starts[i], uy, uw, uh,
         fill=C['white'], border=C['amber'], lbl_size=12, sub_size=10)

# ── 프론트 레인 박스 ──────────────────────────────────────────
f_nodes = [
    ('DocumentsPage', 'TC레벨 선택'),
    ('TreeViewPage', '트리 편집 UI'),
    ('TCReviewPage', '검토·승인 UI'),
    ('Excel 다운로드', 'api.exportUrl'),
]
fw = Inches(1.7); fh = Inches(1.3); fy = Inches(2.75)
fx_starts = [Inches(1.6), Inches(4.25), Inches(7.45), Inches(10.95)]
for i, (lbl, sub) in enumerate(f_nodes):
    node(sl, lbl, sub, fx_starts[i], fy, fw, fh,
         fill=C['blue_pale'], border=C['blue_mid'], lbl_size=12, sub_size=10)

# ── 백엔드 레인 박스 ──────────────────────────────────────────
b_nodes = [
    ('② AI 메뉴트리\n추출', 'TREE_PROMPT\n4단계 계층'),
    ('④ TC 생성\n파이프라인', 'RAG + Gemini\n기능별 순차처리'),
    ('진행률\nDB 업데이트', 'progress_current\nprogress_total'),
    ('TC 저장 +\nTC이력 벡터화', 'SQLite\ntc_history.pkl'),
]
bw2 = Inches(1.75); bh2 = Inches(1.65); by2 = Inches(4.75)
bx_starts = [Inches(1.55), Inches(4.2), Inches(7.1), Inches(10.0)]
b_fills  = [C['blue_mid'], C['blue_dark'], C['teal'], C['green']]
for i, (lbl, sub) in enumerate(b_nodes):
    node(sl, lbl, sub, bx_starts[i], by2, bw2, bh2,
         fill=b_fills[i], border=b_fills[i], lbl_size=12, sub_size=10)
    # 텍스트 컬러 흰색으로
    # (node 함수가 fill 기반이라 별도 덮어쓰기)
    txt(sl, lbl, bx_starts[i]+Inches(0.1), by2+Inches(0.2),
        bw2-Inches(0.2), Inches(0.55),
        size=12, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    txt(sl, sub, bx_starts[i]+Inches(0.1), by2+Inches(0.78),
        bw2-Inches(0.2), Inches(0.6),
        size=10, color=C['blue_light'], align=PP_ALIGN.CENTER)

# ── 세로 화살표 (사용자 → 프론트) ───────────────────────────
v_pairs = [(Inches(2.48), Inches(2.25), Inches(2.75)),   # upload→documents
           (Inches(5.1),  Inches(2.25), Inches(2.75)),   # 검토→tree
           (Inches(8.3),  Inches(2.25), Inches(2.75)),   # tc검토→review
           (Inches(11.73),Inches(2.25), Inches(2.75)),   # 다운로드
           ]
for x, y1, y2 in v_pairs:
    arrow_v(sl, x, y1, y2, color=C['blue_mid'])

# ── 세로 화살표 (프론트 → 백엔드) ───────────────────────────
v2_pairs = [(Inches(2.45), Inches(4.05), Inches(4.75)),
            (Inches(5.1),  Inches(4.05), Inches(4.75)),
            (Inches(8.3),  Inches(4.05), Inches(4.75)),
            (Inches(11.73),Inches(4.05), Inches(4.75)),
            ]
for x, y1, y2 in v2_pairs:
    arrow_v(sl, x, y1, y2, color=C['green'])

# ── 가로 화살표 (백엔드 내부 흐름) ─────────────────────────
arrow_h(sl, Inches(3.3), Inches(5.58), Inches(4.2), color=C['white'])
arrow_h(sl, Inches(5.95), Inches(5.58), Inches(7.1), color=C['white'])
arrow_h(sl, Inches(8.85), Inches(5.58), Inches(10.0), color=C['white'])

# ── 폴링 표시 ─────────────────────────────────────────────
rect(sl, Inches(7.1), Inches(6.55), Inches(5.7), Inches(0.5),
     fill=C['blue_light'], line_color=C['blue_mid'])
txt(sl, '3초 폴링: status + progress_current / progress_total → 프론트 진행률 바 업데이트',
    Inches(7.2), Inches(6.62), Inches(5.5), Inches(0.38),
    size=11, color=C['blue_dark'])

prs1.save('docs/Aegis_시스템동작_다이어그램.pptx')
print('saved: docs/Aegis_시스템동작_다이어그램.pptx')


# ════════════════════════════════════════════════════════════════════
#  PPT 2 — RAG 구조 다이어그램
# ════════════════════════════════════════════════════════════════════
prs2 = Presentation()
prs2.slide_width = W; prs2.slide_height = H
blank2 = prs2.slide_layouts[6]

sl2 = prs2.slides.add_slide(blank2)
rect(sl2, 0, 0, W, H, fill=C['gray_light'])

# 헤더
rect(sl2, 0, 0, W, Inches(0.85), fill=C['blue_dark'])
txt(sl2, 'Aegis QA Assistant — RAG 구조 다이어그램',
    Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
    size=22, bold=True, color=C['white'])
txt(sl2, '매뉴얼 · 결함이력 · TC이력 벡터화 → 유사도 검색 → TC 생성 프롬프트에 자동 주입',
    Inches(0.4), Inches(0.62), Inches(12), Inches(0.28),
    size=12, color=C['blue_light'])

# ── 좌측: 데이터 소스 ─────────────────────────────────────
src_title = [
    ('📖  매뉴얼 PDF', 'XpERP 8종 (검침·부과·수납·\n회계·인사급여·입주자 등)\n총 362+ 청크'),
    ('🐛  테스트 결과 Excel', 'QA 수행 후 업로드\n검증결과=Fail 행 자동 추출\n결함내용·재현절차·Jira'),
    ('📝  생성된 TC', 'TC 생성/재생성 완료 시\n자동 누적 저장\n문서별 갱신'),
]
src_colors = [C['blue_mid'], C['red'], C['green']]
src_y = [Inches(1.05), Inches(3.0), Inches(5.0)]
src_w = Inches(2.8); src_h = Inches(1.6)

for i, ((title, desc), col) in enumerate(zip(src_title, src_colors)):
    rect(sl2, Inches(0.25), src_y[i], src_w, src_h,
         fill=col, line_color=col)
    txt(sl2, title,
        Inches(0.35), src_y[i]+Inches(0.08), src_w-Inches(0.2), Inches(0.38),
        size=12, bold=True, color=C['white'])
    txt(sl2, desc,
        Inches(0.35), src_y[i]+Inches(0.48), src_w-Inches(0.2), Inches(1.0),
        size=11, color=C['blue_light'] if col == C['blue_mid'] else C['white'])

# ── 인제스트 화살표 ──────────────────────────────────────
ing_y_mid = [src_y[i]+Inches(0.8) for i in range(3)]
for y in ing_y_mid:
    arrow_h(sl2, Inches(3.05), y, Inches(3.7), color=C['gray_mid'])

# ── 인제스트 박스 ─────────────────────────────────────────
proc_info = [
    ('벡터화 파이프라인', '800자 청크 분할\nGemini Embedding API\n배치 50개 처리'),
    ('벡터화 파이프라인', 'Fail 행 추출\nGemini Embedding API\n기능영역/결함/절차 포맷'),
    ('자동 벡터화', 'TC 생성 완료 시 자동 호출\nGemini Embedding API\n문서별 교체 저장'),
]
proc_colors = [C['blue_light'], C['red_light'], C['green_light']]
proc_border = [C['blue_mid'], C['red'], C['green']]
proc_w = Inches(2.4)

for i in range(3):
    rect(sl2, Inches(3.7), src_y[i], proc_w, src_h,
         fill=proc_colors[i], line_color=proc_border[i], line_width=1.5)
    txt(sl2, proc_info[i][0],
        Inches(3.8), src_y[i]+Inches(0.08), proc_w-Inches(0.2), Inches(0.35),
        size=11, bold=True, color=proc_border[i])
    txt(sl2, proc_info[i][1],
        Inches(3.8), src_y[i]+Inches(0.46), proc_w-Inches(0.2), Inches(1.0),
        size=10.5, color=C['gray_mid'])

# ── 저장소 화살표 ─────────────────────────────────────────
for y in ing_y_mid:
    arrow_h(sl2, Inches(6.1), y, Inches(6.75), color=C['gray_mid'])

# ── 벡터 DB ───────────────────────────────────────────────
db_info = [
    ('manual_xperp.pkl', 'SimpleVectorStore\n362+ 청크'),
    ('defect_history.pkl', 'SimpleVectorStore\nFail TC 누적'),
    ('tc_history.pkl', 'SimpleVectorStore\n자동 갱신'),
]
db_colors = [C['blue_mid'], C['red'], C['green']]
db_w = Inches(2.1)

for i, ((fname, desc), col) in enumerate(zip(db_info, db_colors)):
    rect(sl2, Inches(6.75), src_y[i], db_w, src_h,
         fill=C['white'], line_color=col, line_width=2.5)
    # DB 아이콘
    rect(sl2, Inches(6.75), src_y[i], db_w, Inches(0.38), fill=col)
    txt(sl2, '🗄  ' + fname,
        Inches(6.85), src_y[i]+Inches(0.04), db_w-Inches(0.2), Inches(0.32),
        size=10, bold=True, color=C['white'])
    txt(sl2, desc,
        Inches(6.85), src_y[i]+Inches(0.5), db_w-Inches(0.2), Inches(0.8),
        size=10.5, color=C['gray_mid'])

# ── 검색 화살표 ───────────────────────────────────────────
for y in ing_y_mid:
    arrow_h(sl2, Inches(8.85), y, Inches(9.5), color=C['blue_mid'])
    txt(sl2, '유사도\n검색', Inches(8.88), y-Inches(0.35),
        Inches(0.65), Inches(0.38), size=9, color=C['blue_mid'])

# ── 검색 결과 통합 ────────────────────────────────────────
rect(sl2, Inches(9.5), Inches(2.6), Inches(2.3), Inches(1.6),
     fill=C['blue_pale'], line_color=C['blue_mid'], line_width=1.5)
txt(sl2, 'RAG 컨텍스트\n통합',
    Inches(9.6), Inches(2.68), Inches(2.1), Inches(0.42),
    size=12, bold=True, color=C['blue_dark'], align=PP_ALIGN.CENTER)
items = ['매뉴얼 상위 8개', '결함이력 상위 5개', 'TC이력 상위 5개']
for j, it in enumerate(items):
    txt(sl2, '• ' + it,
        Inches(9.6), Inches(3.12)+j*Inches(0.28),
        Inches(2.1), Inches(0.28), size=10.5, color=C['gray_mid'])

# 통합→프롬프트 화살표
arrow_h(sl2, Inches(11.8), Inches(3.4), Inches(12.15), color=C['blue_mid'])

# ── Gemini API 박스 ───────────────────────────────────────
rect(sl2, Inches(12.15), Inches(2.2), Inches(0.9), Inches(2.42),
     fill=C['blue_dark'], line_color=C['blue_dark'])
txt(sl2, 'G\ne\nm\ni\nn\ni\nAPI',
    Inches(12.2), Inches(2.35), Inches(0.8), Inches(2.1),
    size=11, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

# ── 쿼리 입력 화살표 (기능 카테고리 → 검색) ──────────────
rect(sl2, Inches(9.5), Inches(1.05), Inches(2.3), Inches(1.2),
     fill=C['amber_light'], line_color=C['amber'], line_width=1.5)
txt(sl2, '기능 카테고리\n(TC 생성 대상)',
    Inches(9.6), Inches(1.15), Inches(2.1), Inches(0.45),
    size=11, bold=True, color=C['amber'], align=PP_ALIGN.CENTER)
txt(sl2, '리프노드 경로\n예: 부과 > 부과화면 > 입력폼',
    Inches(9.6), Inches(1.62), Inches(2.1), Inches(0.5),
    size=10, color=C['gray_mid'], align=PP_ALIGN.CENTER)

# 카테고리 → 각 DB 검색 화살표
for y in ing_y_mid:
    arrow_v(sl2, Inches(10.65), Inches(2.25), y, color=C['amber'])

# ── 하단 설명 ─────────────────────────────────────────────
rect(sl2, Inches(0.25), Inches(6.72), Inches(12.8), Inches(0.58),
     fill=C['blue_dark'])
notes = [
    'SimpleVectorStore: ChromaDB 대신 numpy + pickle로 자체 구현 (MSVC 미설치 환경 호환)',
    '|  임계값: 매뉴얼 0.50 / 결함이력 0.45 / TC이력 0.45  |  임베딩: models/gemini-embedding-2',
]
txt(sl2, '  '.join(notes),
    Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.42),
    size=11, color=C['blue_light'])

# ── 수동 vs 자동 배지 ────────────────────────────────────
badge(sl2, '수동 업로드', Inches(0.25), src_y[0]+src_h+Inches(0.08),
      Inches(1.1), Inches(0.28), fill=C['blue_mid'])
badge(sl2, '수동 업로드', Inches(0.25), src_y[1]+src_h+Inches(0.08),
      Inches(1.1), Inches(0.28), fill=C['red'])
badge(sl2, '자동 누적', Inches(0.25), src_y[2]+src_h+Inches(0.08),
      Inches(1.1), Inches(0.28), fill=C['green'])

prs2.save('docs/Aegis_RAG_구조_다이어그램.pptx')
print('saved: docs/Aegis_RAG_구조_다이어그램.pptx')
