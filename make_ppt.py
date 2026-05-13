from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── 색상 팔레트 ──────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x1e, 0x40, 0xaf)   # #1e40af
BLUE_MID    = RGBColor(0x25, 0x63, 0xeb)   # #2563eb
BLUE_LIGHT  = RGBColor(0xdb, 0xea, 0xfe)   # #dbeafe
GRAY_DARK   = RGBColor(0x11, 0x18, 0x27)   # #111827
GRAY_MID    = RGBColor(0x37, 0x41, 0x51)   # #374151
GRAY_LIGHT  = RGBColor(0xf3, 0xf4, 0xf6)   # #f3f4f6
WHITE       = RGBColor(0xff, 0xff, 0xff)
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
AMBER       = RGBColor(0xd9, 0x77, 0x06)

W = Inches(13.33)   # 와이드 슬라이드 너비
H = Inches(7.5)     # 와이드 슬라이드 높이

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 헬퍼 함수 ────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def card(slide, x, y, w, h, title, lines, title_color=BLUE_MID,
         bg=WHITE, line_color=BLUE_LIGHT, title_size=15, body_size=12.5):
    add_rect(slide, x, y, w, h, fill=bg, line=line_color)
    add_rect(slide, x, y, w, Inches(0.38), fill=title_color)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.05),
             w - Inches(0.3), Inches(0.32),
             size=title_size, bold=True, color=WHITE)
    txb = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.42),
        w - Inches(0.3), h - Inches(0.55))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(body_size)
        run.font.color.rgb = GRAY_MID


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.1), fill=BLUE_DARK)
    add_text(slide, title, Inches(0.5), Inches(0.12), Inches(11), Inches(0.55),
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.65), Inches(11), Inches(0.38),
                 size=14, color=BLUE_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=BLUE_DARK)
add_rect(sl, 0, Inches(2.8), W, Inches(1.9), fill=BLUE_MID)

add_text(sl, '🛡 Aegis QA Assistant',
         Inches(1), Inches(1.0), Inches(11), Inches(0.9),
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 'AI 기반 테스트케이스 자동 생성 시스템',
         Inches(1), Inches(1.95), Inches(11), Inches(0.6),
         size=22, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, 'PDF 기획서  →  메뉴트리 추출  →  TC 자동 생성  →  Excel 리포트',
         Inches(0.8), Inches(2.92), Inches(11.5), Inches(0.55),
         size=17, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, '개발팀 공유용  |  2026.05',
         Inches(1), Inches(6.6), Inches(11), Inches(0.5),
         size=13, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 프로젝트 목적
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '프로젝트 목적', 'QA 단계의 인력·비용·시간 절감 + TC 품질 자동 고도화')

pain = [
    '기획 완료 후 TC 수작업 작성에 평균 2~3일 소요',
    'QA 엔지니어 역량에 따라 TC 품질 편차 발생',
    '과거 결함 이력이 다음 프로젝트에 반영되지 않음',
    '기획서 변경 시 TC 전면 재작성 부담',
]
sol = [
    'AI가 PDF 기획서를 분석, TC 초안을 수 분 내 자동 생성',
    '메뉴트리 추출로 테스트 범위를 사전에 시각화·확정',
    '매뉴얼·결함이력·TC이력 RAG로 품질 자동 고도화',
    'TC 생성 레벨(1~5) 조절로 일정·품질 균형 확보',
]
card(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(5.6),
     '😓  기존 Pain Point', pain, title_color=RGBColor(0xdc,0x26,0x26),
     line_color=RGBColor(0xfe,0xe2,0xe2), body_size=13.5)
card(sl, Inches(6.7), Inches(1.25), Inches(6.2), Inches(5.6),
     '✅  Aegis QA로 해결', sol, title_color=GREEN,
     line_color=RGBColor(0xdc,0xfc,0xe7), body_size=13.5)


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 개발 방식
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '개발 방식', 'AI 페어 프로그래밍 — Claude Code와 협업 개발')

add_rect(sl, Inches(0.4), Inches(1.25), W - Inches(0.8), Inches(1.35), fill=WHITE, line=BLUE_LIGHT)
txb = sl.shapes.add_textbox(Inches(0.7), Inches(1.35), W - Inches(1.4), Inches(1.1))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, '개발자가 요구사항과 방향을 제시 → Claude Code(AI)가 코드 작성·수정 → 개발자가 검토·확정하는 반복 사이클',
         size=14, color=GRAY_MID)
add_para(tf, '→  3~4주 내 MVP에서 운영 수준 시스템 완성  |  백엔드(FastAPI) + 프론트(React) + AI 파이프라인 전 영역 커버',
         size=13, color=BLUE_MID, bold=True, space_before=4)

phases = [
    ('Phase 1', 'FastAPI 백엔드\n+ TC 생성 엔진'),
    ('Phase 2', 'RAG 시스템\n(3종 벡터 DB)'),
    ('Phase 3', 'React\nTC 검토 UI'),
    ('Phase 4', '변경유형\n자동 감지'),
    ('Phase 4.5', 'TC 생성\n레벨 시스템'),
    ('Phase 4.6', '메뉴트리 플로우\n+ QA 룰셋'),
    ('Phase 4.7', 'AI 프롬프트\n도메인 특화'),
]
bw = Inches(1.6); gap = Inches(0.2)
start_x = Inches(0.4)
for i, (ph, desc) in enumerate(phases):
    x = start_x + i * (bw + gap)
    add_rect(sl, x, Inches(2.75), bw, Inches(3.6), fill=BLUE_DARK if i >= 4 else BLUE_MID)
    add_text(sl, ph, x, Inches(2.85), bw, Inches(0.45),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x, Inches(3.35), bw, Inches(1.5),
             size=11, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    if i >= 4:
        add_text(sl, '최근', x, Inches(4.9), bw, Inches(0.35),
                 size=10, color=RGBColor(0xfb,0xbf,0x24), align=PP_ALIGN.CENTER, bold=True)

add_text(sl, '✅ 완료', Inches(0.4), Inches(6.45), Inches(4), Inches(0.4),
         size=12, color=GREEN)
add_text(sl, '⭐ 최근 완료 (현재 버전)', Inches(4.5), Inches(6.45), Inches(5), Inches(0.4),
         size=12, color=AMBER)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 시스템 플로우
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '시스템 동작 플로우', '사용자 액션 → AI 분석 → 메뉴트리 검토 → TC 생성 → 검토 → Excel')

steps = [
    ('①', 'PDF\n업로드', BLUE_MID),
    ('②', '메뉴트리\nAI 추출', BLUE_MID),
    ('③', '사용자\n검토·편집', GREEN),
    ('④', 'TC\n생성 시작', BLUE_MID),
    ('⑤', 'RAG+AI\nTC 생성', BLUE_DARK),
    ('⑥', 'TC\n검토·승인', GREEN),
    ('⑦', 'Excel\n다운로드', RGBColor(0x92,0x40,0x0e)),
]
bw = Inches(1.5); by = Inches(1.6); bh = Inches(2.2)
gap = Inches(0.28)
for i, (num, label, col) in enumerate(steps):
    x = Inches(0.25) + i * (bw + gap)
    add_rect(sl, x, by, bw, bh, fill=col)
    add_text(sl, num, x, by + Inches(0.15), bw, Inches(0.5),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, by + Inches(0.6), bw, Inches(1.4),
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + bw + Inches(0.05)
        add_text(sl, '▶', ax, by + Inches(0.85), Inches(0.22), Inches(0.4),
                 size=13, color=GRAY_MID, align=PP_ALIGN.CENTER)

# 하단 AI 박스
add_rect(sl, Inches(0.25), Inches(4.05), Inches(12.8), Inches(1.4), fill=WHITE, line=BLUE_LIGHT)
add_text(sl, '⚡  AI 엔진 (Google Gemini API)',
         Inches(0.45), Inches(4.1), Inches(3.5), Inches(0.4),
         size=13, bold=True, color=BLUE_DARK)
labels = [
    'TREE_PROMPT  →  기획서 분석 + 4단계 메뉴트리 추출',
    'TC_GENERATE_PROMPT  →  기능별 TC 생성 (변경유형·레벨·RAG 반영)',
    'ANALYZE_PROMPT  →  기능 목록 추출 + 변경유형 자동 판단',
]
for j, lb in enumerate(labels):
    add_text(sl, '• ' + lb, Inches(0.45), Inches(4.55) + j * Inches(0.28),
             Inches(12.2), Inches(0.28), size=11.5, color=GRAY_MID)

# 진행률 표시 안내
add_rect(sl, Inches(0.25), Inches(5.6), Inches(12.8), Inches(0.7), fill=BLUE_LIGHT)
add_text(sl, '📊  TC 생성 중 실시간 진행률 표시  :  기능 N/M 처리 중  ·  경과 시간  ·  예상 완료까지 N분',
         Inches(0.5), Inches(5.72), Inches(12), Inches(0.45),
         size=13, color=BLUE_DARK, bold=True)

add_text(sl, '🕐  KST 기준 날짜/시간 표시',
         Inches(0.25), Inches(6.45), Inches(5), Inches(0.4),
         size=11, color=GRAY_MID)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — RAG 시스템
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, 'RAG 시스템 — 품질 자동 고도화', '매뉴얼·결함이력·TC이력을 벡터화하여 TC 생성 시 자동 반영')

sources = [
    ('📖  매뉴얼 벡터 DB', 'manual_xperp.pkl', [
        'XpERP 업무 매뉴얼 PDF 8종',
        '총 362+ 청크 (800자 단위)',
        '검침·부과·수납·회계·인사급여 등',
        '→ 실제 스펙·입력조건·오류메시지 반영',
    ], BLUE_MID),
    ('🐛  결함이력 벡터 DB', 'defect_history.pkl', [
        '과거 테스트 결과 Excel 업로드',
        'Fail TC 자동 추출·벡터화',
        '유사 기능의 비정상/경계값 강화',
        '→ 반복 결함 사전 차단',
    ], RGBColor(0xdc,0x26,0x26)),
    ('📝  TC이력 벡터 DB', 'tc_history.pkl', [
        'TC 생성/재생성 시 자동 누적',
        '문서별 TC 구조·패턴 색인',
        '유사 기능 TC 스타일 참고',
        '→ 사이클 반복 시 품질 향상',
    ], GREEN),
]
cw = Inches(3.8)
for i, (title, fname, lines, col) in enumerate(sources):
    x = Inches(0.4) + i * (cw + Inches(0.35))
    card(sl, x, Inches(1.3), cw, Inches(3.5), title, lines,
         title_color=col, body_size=13)
    add_text(sl, fname, x, Inches(4.9), cw, Inches(0.35),
             size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.65), fill=WHITE, line=BLUE_LIGHT)
add_text(sl, '🔍  유사도 검색 (코사인 유사도)  →  상위 결과를 TC 생성 프롬프트에 자동 주입  →  Gemini API가 스펙·결함·패턴 반영하여 TC 생성',
         Inches(0.6), Inches(5.25), Inches(12), Inches(0.45),
         size=13, color=BLUE_DARK)

add_rect(sl, Inches(0.4), Inches(5.95), Inches(12.5), Inches(1.15), fill=BLUE_DARK)
add_text(sl, '💡  SimpleVectorStore  —  자체 구현 벡터 DB (numpy 기반)',
         Inches(0.65), Inches(6.0), Inches(12), Inches(0.4),
         size=13, bold=True, color=WHITE)
add_text(sl, 'ChromaDB(C++ 컴파일 필요) 대신 순수 Python + numpy로 구현  ·  pickle 파일로 로컬 저장  ·  코사인 유사도 검색',
         Inches(0.65), Inches(6.42), Inches(12), Inches(0.55),
         size=12, color=BLUE_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 핵심 기능
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '핵심 기능', '메뉴트리 · TC 레벨 · QA 룰셋 · 도메인 특화 프롬프트')

features = [
    ('🌲  메뉴트리 추출', [
        '4단계 계층: 모듈 > 화면 > 영역 > 기능',
        '리프 노드 = TC 생성 단위',
        '사용자가 직접 편집·제외 설정 가능',
        'Excel 다운로드 지원',
    ]),
    ('📊  TC 생성 레벨 (1~5)', [
        'Lv1 핵심검증 · Lv2 표준 · Lv3 정밀(기본)',
        'Lv4 심층 · Lv5 전수검증',
        '레벨별 기능당 생성 수: 2~35개',
        '모든 레벨에서 4가지 유형 분포 유지',
    ]),
    ('⚙️  QA 룰셋 관리', [
        '메뉴트리 추출 지침 + TC 생성 지침',
        '프로젝트별 룰셋 선택 적용',
        '시스템 기본 룰셋 + 사용자 정의 룰셋',
        'UI에서 직접 생성·수정·삭제',
    ]),
    ('🤖  도메인 특화 프롬프트', [
        '아파트·공동주택 관리 ERP 전문가 역할',
        '회계/부과/검침/인사급여/입주관리 도메인',
        '공동주택관리법·회계처리기준 규정 반영',
        '경계값·동등분할·탐색적 테스트 기법 명시',
    ]),
]
cw = Inches(2.9); ch = Inches(4.9)
for i, (title, lines) in enumerate(features):
    x = Inches(0.4) + i * (cw + Inches(0.5))
    card(sl, x, Inches(1.3), cw, ch, title, lines, body_size=13)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 기술 스택
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '기술 스택', 'Python + FastAPI + React + Google Gemini API')

stacks = [
    ('🖥️  백엔드', [
        'Python 3.12  +  FastAPI',
        'SQLAlchemy ORM  /  SQLite → PostgreSQL',
        'Uvicorn (ASGI 서버)',
        'python-multipart, openpyxl, pypdf',
    ], BLUE_MID),
    ('⚛️  프론트엔드', [
        'React 18  +  Vite',
        'React Router DOM 6',
        'Axios (HTTP 클라이언트)',
        '전역 CSS (index.css)',
    ], RGBColor(0x06,0xb6,0xd4)),
    ('🤖  AI / 벡터', [
        'Google Gemini API  (google-genai)',
        'gemini-2.5-flash  /  flash-lite',
        'Gemini Embedding API',
        'SimpleVectorStore (numpy)',
    ], BLUE_DARK),
    ('🗄️  인프라', [
        'SQLite (개발)  →  PostgreSQL (예정)',
        'FastAPI 정적파일 서빙 (운영 단일 서버)',
        'GitHub 형상관리',
        'Windows 10 Pro / VS Code',
    ], GRAY_MID),
]
cw = Inches(2.9)
for i, (title, lines, col) in enumerate(stacks):
    x = Inches(0.4) + i * (cw + Inches(0.47))
    card(sl, x, Inches(1.3), cw, Inches(4.2), title, lines,
         title_color=col, body_size=13)

add_rect(sl, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.55), fill=WHITE, line=BLUE_LIGHT)
add_text(sl, '⚠️  Python 3.12 필수 (3.14 pydantic-core 휠 미제공)  ·  pip 설치 시 사내 SSL 우회 옵션 필요  ·  MSVC 불필요 (순수 Python 패키지만 사용)',
         Inches(0.6), Inches(5.79), Inches(12), Inches(0.38),
         size=11.5, color=GRAY_MID)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 향후 계획
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, '향후 계획', 'Phase 5~8 로드맵')

roadmap = [
    ('Phase 5', '프론트+백 통합 배포', '단일 서버로 프론트/백 동시 서빙\nnpm build → FastAPI 정적파일 서빙', BLUE_MID, '단기'),
    ('Phase 6', 'TC 수동 결과 기록', '테스트 수행 결과를 시스템 내 기록\n결함 이력 자동 누적 강화', BLUE_MID, '단기'),
    ('Phase 7', 'Playwright 자동화', '웹 크롤러로 UI 수집\nPlaywright 자동화 TC 생성', BLUE_DARK, '중기'),
    ('Phase 8', '팀 협업 + 스케일아웃', '멀티유저·권한 관리\nCelery + Redis 워커 확장', GRAY_MID, '장기'),
]
bw = Inches(2.9); by = Inches(1.35)
for i, (ph, title, desc, col, term) in enumerate(roadmap):
    x = Inches(0.4) + i * (bw + Inches(0.44))
    add_rect(sl, x, by, bw, Inches(4.7), fill=col)
    add_text(sl, ph, x, by + Inches(0.1), bw, Inches(0.4),
             size=13, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(sl, title, x, by + Inches(0.5), bw, Inches(0.55),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + Inches(0.8), by + Inches(1.1), bw - Inches(1.6), Inches(0.32),
             fill=RGBColor(0xff,0xff,0xff) if term=='단기' else BLUE_LIGHT)
    add_text(sl, term, x + Inches(0.8), by + Inches(1.12), bw - Inches(1.6), Inches(0.28),
             size=11, bold=True,
             color=BLUE_DARK if term=='단기' else GRAY_MID,
             align=PP_ALIGN.CENTER)
    add_text(sl, desc, x + Inches(0.15), by + Inches(1.55),
             bw - Inches(0.3), Inches(2.8),
             size=13, color=BLUE_LIGHT)

add_rect(sl, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.65), fill=WHITE, line=BLUE_LIGHT)
add_text(sl, '💬  문의 및 개선 요청  :  iterp@aegisep.com  ·  GitHub: handh228-star/aegis_qa_assistant',
         Inches(0.65), Inches(6.3), Inches(12), Inches(0.45),
         size=13, color=GRAY_MID)


# ── 저장 ────────────────────────────────────────────────────
out = 'docs/Aegis_QA_소개_개발팀공유.pptx'
prs.save(out)
print('saved: ' + out)
